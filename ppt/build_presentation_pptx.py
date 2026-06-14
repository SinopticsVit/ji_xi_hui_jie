#!/usr/bin/env python3
"""Build the capabilities PPTX from the spec markdown files.

Generates maretials/presentation-ru.pptx (28 slides, 16:9):
  - slides 1-6: branded Jixi Huijie / CASC deck (additonal pages/
    additional-slides-spec-ru.md), with photos from video/;
  - slides 7-28: the anonymized capabilities deck (presentation-spec-ru.md).

Slide content is declared as Python data structures below (taken 1:1 from the
spec tables). Images are cropped/masked with Pillow and validated against a
whitelist so no catalog-items thumbnails, supplier branding, blacklisted logo
or CJK headers leak in.

Usage:
  cd aero_secondhend/ppt
  pip install -r requirements.txt
  python build_presentation_pptx.py              # -> maretials/presentation-ru.pptx
  python build_presentation_pptx.py --keep-temp  # keep _pptx_assets/ temp crops
"""

from __future__ import annotations

import argparse
import shutil
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

HERE = Path(__file__).resolve().parent
MATERIALS_DIR = HERE / "maretials"
VIDEO_DIR = HERE / "video"
ASSETS_DIR = HERE.parent / "assets"
TEMP_DIR = MATERIALS_DIR / "_pptx_assets"
OUTPUT_PPTX = MATERIALS_DIR / "presentation-ru.pptx"

YUHUAN = "20260406-yuhuan-tianrun/images"
QIFENG = "qifeng-2022/images"
VIDEO = "video"

# ---------------------------------------------------------------------------
# Image whitelist (spec section 6). Anything else aborts the build.
# Paths are resolved relative to MATERIALS_DIR by default, or to HERE for the
# branded Jixi Huijie slides (video/* assets).
# ---------------------------------------------------------------------------

WHITELIST = {
    f"{YUHUAN}/slide-10.png",
    f"{YUHUAN}/slide-11.png",
    f"{YUHUAN}/slide-12.png",
    f"{YUHUAN}/slide-13.png",
    f"{YUHUAN}/slide-14.png",
    f"{YUHUAN}/slide-15.png",
    f"{YUHUAN}/slide-16.png",
    f"{QIFENG}/slide-11.png",
    f"{QIFENG}/slide-12.png",
    f"{QIFENG}/slide-13.png",
    f"{QIFENG}/slide-14.png",
    f"{QIFENG}/slide-15.png",
}

# Branded additional slides (Jixi Huijie). video/* live under HERE; the xlsx
# catalog crops live under MATERIALS_DIR like the rest.
ADDITIONAL_WHITELIST = {
    f"{VIDEO}/helocopter1.jpg",
    f"{VIDEO}/helocopter2.jpg",
    f"{VIDEO}/helocopter3.jpg",
    f"{VIDEO}/airplane.jpg",
    f"{VIDEO}/agr_airplane.jpg",
    f"{VIDEO}/production1.png",
    f"{VIDEO}/production2.png",
    f"{VIDEO}/cert.png",
    f"{YUHUAN}/_xlsx_slide-02.png",
    f"{YUHUAN}/_xlsx_slide-03.png",
}

# Explicitly forbidden assets (spec section 5.4).
BLACKLIST = {
    (ASSETS_DIR / "logo-eagle-red.png").resolve(),
}

# ---------------------------------------------------------------------------
# Design tokens (spec section 2)
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
ACCENT = RGBColor(0x00, 0x78, 0xA8)
BG = RGBColor(0xFF, 0xFF, 0xFF)
BG_SOFT = RGBColor(0xF5, 0xF7, 0xFA)
RULE = RGBColor(0xCB, 0xD5, 0xE1)
TEXT = RGBColor(0x1F, 0x2A, 0x44)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xF8, 0xFA, 0xFC)
ACCENT_SOFT = RGBColor(0xE8, 0xF4, 0xF8)

FONT_HEAD = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"

SIDE_M = Inches(0.55)
HEADER_H = Inches(0.50)
TOTAL_SLIDES = 28

# ---------------------------------------------------------------------------
# Low-level shape helpers
# ---------------------------------------------------------------------------


def _solid_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _set_bg(slide, color: RGBColor) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid_fill(rect, color)
    # send to back
    sp = rect._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def _set_text(
    tf,
    text: str,
    *,
    font: str = FONT_BODY,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _new_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_accent_line(slide, left, top, width, *, color: RGBColor = ACCENT,
                     height=Pt(3)) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _solid_fill(bar, color)


def _add_header(slide, title: str) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
    _solid_fill(band, NAVY)
    title_box = slide.shapes.add_textbox(
        SIDE_M, Inches(0.62), SLIDE_W - 2 * SIDE_M, Inches(0.7)
    )
    _set_text(title_box.text_frame, title, font=FONT_HEAD, size=26, bold=True,
              color=NAVY)
    _add_accent_line(slide, SIDE_M, Inches(1.32), Inches(2.4))


def _add_footer(slide, idx: int) -> None:
    box = slide.shapes.add_textbox(
        SLIDE_W - Inches(2.0), SLIDE_H - Inches(0.42), Inches(1.6), Inches(0.3)
    )
    _set_text(box.text_frame, f"{idx} / {TOTAL_SLIDES}", font=FONT_BODY, size=10,
              color=MUTED, align=PP_ALIGN.RIGHT)


def _add_bullets(slide, left, top, width, height, bullets: list[str], *,
                 size: int = 17, color: RGBColor = TEXT, gap: float = 0.18) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(size * gap * 4)
        p.line_spacing = 1.15
        dot = p.add_run()
        dot.text = "•  "
        dot.font.name = FONT_BODY
        dot.font.size = Pt(size)
        dot.font.bold = True
        dot.font.color.rgb = ACCENT
        run = p.add_run()
        run.text = line
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color


def _add_table(slide, left, top, width, header: list[str], rows: list[list[str]],
               *, col_widths: list[float] | None = None, header_size: int = 13,
               body_size: int = 12) -> None:
    n_rows = len(rows) + 1
    n_cols = len(header)
    row_h = Inches(0.42)
    gfx = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                 row_h * n_rows)
    table = gfx.table
    table.first_row = False
    table.horz_banding = False

    if col_widths:
        total = sum(col_widths)
        for c, frac in enumerate(col_widths):
            table.columns[c].width = Emu(int(width * frac / total))

    for c, label in enumerate(header):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(2)
        cell.margin_bottom = Pt(2)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = label
        run.font.name = FONT_HEAD
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = WHITE

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_ALT if r % 2 == 0 else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            run.font.name = FONT_BODY
            run.font.size = Pt(body_size)
            run.font.color.rgb = TEXT


# ---------------------------------------------------------------------------
# Image pipeline (spec section 4)
# ---------------------------------------------------------------------------


def _resolve_source(rel_path: str) -> Path:
    """video/* assets live under HERE, everything else under MATERIALS_DIR."""
    base = HERE if rel_path.startswith(f"{VIDEO}/") else MATERIALS_DIR
    return base / rel_path


def assert_allowed(rel_path: str) -> None:
    assert "catalog-items" not in rel_path, f"blacklist (catalog-items): {rel_path}"
    src = _resolve_source(rel_path).resolve()
    assert src not in BLACKLIST, f"blacklist: {rel_path}"
    allowed = WHITELIST | ADDITIONAL_WHITELIST
    assert rel_path in allowed, f"not in whitelist: {rel_path}"


def prepare_image(rel_path: str, *, crop_px=None, crop_frac=None,
                  top_mask: float = 0.0, bottom_mask: float = 0.0,
                  mask_color=(11, 42, 74)) -> Path:
    """Crop/mask a whitelisted source image (PNG/JPG) and return a temp path.

    crop_px:   (left, top, right, bottom) in source pixels.
    crop_frac: (x, y, w, h) as fractions of width/height.
    top_mask:  fraction of the cropped height to paint over at the top
               (hides CJK headlines burnt into video stills).
    bottom_mask: fraction of the cropped height to paint over at the bottom
                 (hides footers / URLs / branding).
    """
    assert_allowed(rel_path)
    src = _resolve_source(rel_path)
    if not src.is_file():
        raise FileNotFoundError(f"missing image: {src}")

    img = Image.open(src).convert("RGB")
    w, h = img.size

    if crop_px is not None:
        img = img.crop(crop_px)
    elif crop_frac is not None:
        x, y, cw, ch = crop_frac
        box = (int(w * x), int(h * y), int(w * (x + cw)), int(h * (y + ch)))
        img = img.crop(box)

    cw, ch = img.size
    if top_mask > 0:
        band_h = int(ch * top_mask)
        if band_h > 0:
            img.paste(Image.new("RGB", (cw, band_h), mask_color), (0, 0))
    if bottom_mask > 0:
        band_h = int(ch * bottom_mask)
        if band_h > 0:
            img.paste(Image.new("RGB", (cw, band_h), mask_color), (0, ch - band_h))

    TEMP_DIR.mkdir(parents=True, exist_ok=True)
    safe = rel_path.replace("/", "__").rsplit(".", 1)[0]
    out = TEMP_DIR / f"crop__{safe}.png"
    img.save(out, format="PNG")
    return out


def _add_image_fit(slide, img_path: Path, left, top, max_w, max_h, *,
                   border: bool = True, center_h: bool = True):
    with Image.open(img_path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    draw_w = int(iw * scale)
    draw_h = int(ih * scale)
    if center_h:
        left = left + (max_w - draw_w) // 2
    top = top + (max_h - draw_h) // 2
    pic = slide.shapes.add_picture(str(img_path), left, top, draw_w, draw_h)
    if border:
        pic.line.color.rgb = RULE
        pic.line.width = Pt(1)
    return pic


# ---------------------------------------------------------------------------
# Layout renderers
# ---------------------------------------------------------------------------


def render_title(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG_SOFT)
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.6),
                                         SLIDE_W - Inches(2.0), Inches(1.4))
    _set_text(title_box.text_frame, data["title"], font=FONT_HEAD, size=38,
              bold=True, color=NAVY, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    _add_accent_line(slide, (SLIDE_W - Inches(2.2)) // 2, Inches(4.05),
                     Inches(2.2))
    if data.get("subtitle"):
        sub = slide.shapes.add_textbox(Inches(1.5), Inches(4.35),
                                       SLIDE_W - Inches(3.0), Inches(1.6))
        _set_text(sub.text_frame, data["subtitle"], font=FONT_BODY, size=18,
                  color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    if data.get("bullets"):
        _add_bullets(slide, Inches(3.4), Inches(4.5), SLIDE_W - Inches(6.8),
                     Inches(2.4), data["bullets"], size=15, color=TEXT)


def render_section(slide, data: dict, idx: int) -> None:
    _set_bg(slide, NAVY)
    letter = slide.shapes.add_textbox(SIDE_M, Inches(2.0), Inches(3.0),
                                      Inches(2.2))
    _set_text(letter.text_frame, data["letter"], font=FONT_HEAD, size=120,
              bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    title = slide.shapes.add_textbox(Inches(3.0), Inches(2.9),
                                     SLIDE_W - Inches(3.6), Inches(1.4))
    _set_text(title.text_frame, data["title"], font=FONT_HEAD, size=34,
              bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if data.get("note"):
        note = slide.shapes.add_textbox(Inches(3.05), Inches(4.2),
                                        SLIDE_W - Inches(3.6), Inches(1.0))
        _set_text(note.text_frame, data["note"], font=FONT_BODY, size=16,
                  color=ACCENT_SOFT)


def render_split_text_image(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG)
    _add_header(slide, data["title"])
    _add_bullets(slide, SIDE_M, Inches(1.7), Inches(5.4), Inches(4.8),
                 data["bullets"], size=17)

    img_left = Inches(6.4)
    img_top = Inches(1.7)
    img_w = SLIDE_W - img_left - SIDE_M
    img_h = Inches(4.6)
    if data.get("image"):
        path = prepare_image(data["image"], **data.get("image_args", {}))
        _add_image_fit(slide, path, img_left, img_top, img_w, img_h)
        if data.get("caption"):
            cap = slide.shapes.add_textbox(img_left, Inches(6.45), img_w,
                                           Inches(0.5))
            _set_text(cap.text_frame, data["caption"], font=FONT_BODY, size=10,
                      color=MUTED, align=PP_ALIGN.CENTER)
    _add_footer(slide, idx)


def render_image_hero(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG)
    _add_header(slide, data["title"])
    path = prepare_image(data["image"], **data.get("image_args", {}))
    img_left = SIDE_M
    img_top = Inches(1.65)
    img_w = SLIDE_W - 2 * SIDE_M
    img_h = Inches(4.55)
    _add_image_fit(slide, path, img_left, img_top, img_w, img_h)
    if data.get("caption"):
        cap = slide.shapes.add_textbox(SIDE_M, Inches(6.45),
                                       SLIDE_W - 2 * SIDE_M, Inches(0.5))
        _set_text(cap.text_frame, data["caption"], font=FONT_BODY, size=12,
                  color=MUTED, align=PP_ALIGN.CENTER)
    _add_footer(slide, idx)


def render_dual_image(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG)
    _add_header(slide, data["title"])

    content_w = SLIDE_W - 2 * SIDE_M
    gap = Inches(0.4)
    cell_w = (content_w - gap) // 2
    cell_h = Inches(3.9)
    top = Inches(1.65)

    for i, item in enumerate(data["images"]):
        left = SIDE_M + i * (cell_w + gap)
        path = prepare_image(item["image"], **item.get("image_args", {}))
        _add_image_fit(slide, path, left, top, cell_w, cell_h)
        cap = slide.shapes.add_textbox(left, top + cell_h + Inches(0.05),
                                       cell_w, Inches(0.4))
        _set_text(cap.text_frame, item["caption"], font=FONT_HEAD, size=13,
                  bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    if data.get("caption"):
        note = slide.shapes.add_textbox(SIDE_M, Inches(6.35),
                                        SLIDE_W - 2 * SIDE_M, Inches(0.7))
        _set_text(note.text_frame, data["caption"], font=FONT_BODY, size=13,
                  color=MUTED, align=PP_ALIGN.CENTER)
    _add_footer(slide, idx)


def render_matrix(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG)
    _add_header(slide, data["title"])

    if data.get("cards"):
        _render_cards(slide, data["cards"])
    else:
        _add_table(
            slide, SIDE_M, Inches(1.7), SLIDE_W - 2 * SIDE_M,
            data["header"], data["rows"],
            col_widths=data.get("col_widths"),
        )
    if data.get("note"):
        note = slide.shapes.add_textbox(SIDE_M, SLIDE_H - Inches(0.9),
                                        SLIDE_W - 2 * SIDE_M, Inches(0.5))
        _set_text(note.text_frame, data["note"], font=FONT_BODY, size=12,
                  color=MUTED)
    _add_footer(slide, idx)


def _render_cards(slide, cards: list[dict]) -> None:
    content_w = SLIDE_W - 2 * SIDE_M
    gap = Inches(0.5)
    cell_w = (content_w - gap) // 2
    top = Inches(1.8)
    cell_h = Inches(4.2)
    for i, card in enumerate(cards):
        left = SIDE_M + i * (cell_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                     cell_w, cell_h)
        box.fill.solid()
        box.fill.fore_color.rgb = BG_SOFT
        box.line.color.rgb = RULE
        box.line.width = Pt(1)
        head = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.25),
                                        cell_w - Inches(0.6), Inches(0.7))
        _set_text(head.text_frame, card["title"], font=FONT_HEAD, size=20,
                  bold=True, color=NAVY)
        _add_bullets(slide, left + Inches(0.3), top + Inches(1.1),
                     cell_w - Inches(0.6), cell_h - Inches(1.4),
                     card["bullets"], size=15)


def _add_multiline(slide, left, top, width, height, lines: list[str], *,
                   font: str = FONT_BODY, size: int = 12,
                   color: RGBColor = WHITE, align: PP_ALIGN = PP_ALIGN.CENTER,
                   anchor: MSO_ANCHOR = MSO_ANCHOR.TOP) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.1
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color


def _image_columns(slide, items: list[dict], top, cell_h, *,
                   card_bg: bool = False) -> None:
    content_w = SLIDE_W - 2 * SIDE_M
    n = len(items)
    gap = Inches(0.35)
    cell_w = (content_w - gap * (n - 1)) // n
    for i, item in enumerate(items):
        left = SIDE_M + i * (cell_w + gap)
        img_top = top
        img_h = cell_h
        if card_bg:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                         cell_w, cell_h + Inches(0.5))
            box.fill.solid()
            box.fill.fore_color.rgb = BG_SOFT
            box.line.color.rgb = RULE
            box.line.width = Pt(1)
            img_top = top + Inches(0.15)
            img_h = cell_h - Inches(0.2)
        path = prepare_image(item["image"], **item.get("image_args", {}))
        _add_image_fit(slide, path, left + Inches(0.12), img_top,
                       cell_w - Inches(0.24), img_h, border=not card_bg)
        if item.get("caption"):
            cap = slide.shapes.add_textbox(left + Inches(0.12),
                                           top + cell_h + Inches(0.12),
                                           cell_w - Inches(0.24), Inches(0.6))
            _set_text(cap.text_frame, item["caption"], font=FONT_BODY, size=12,
                      color=NAVY, align=PP_ALIGN.CENTER)


def render_title_branded(slide, data: dict, idx: int) -> None:
    _set_bg(slide, ACCENT)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.45),
                                     SLIDE_W - Inches(1.6), Inches(1.0))
    _set_text(title.text_frame, data["title"], font=FONT_HEAD, size=30,
              bold=True, color=WHITE, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    if data.get("subtitle"):
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.45),
                                       SLIDE_W - Inches(1.6), Inches(0.55))
        _set_text(sub.text_frame, data["subtitle"], font=FONT_BODY, size=16,
                  color=ACCENT_SOFT, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE)
    if data.get("image"):
        path = prepare_image(data["image"], **data.get("image_args", {}))
        _add_image_fit(slide, path, Inches(2.4), Inches(2.15),
                       SLIDE_W - Inches(4.8), Inches(3.6))
    if data.get("contacts"):
        _add_multiline(slide, Inches(0.8), SLIDE_H - Inches(1.3),
                       SLIDE_W - Inches(1.6), Inches(1.1), data["contacts"],
                       size=13, color=WHITE, align=PP_ALIGN.CENTER)


def render_split_gallery(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG)
    _add_header(slide, data["title"])

    images = data.get("images", [])
    grid_left = SIDE_M
    grid_top = Inches(1.7)
    grid_w = Inches(6.2)
    cols = 2
    rows = (len(images) + cols - 1) // cols
    gap = Inches(0.25)
    cell_w = (grid_w - gap) // cols
    cell_h = (Inches(4.8) - gap * (rows - 1)) // rows
    for i, item in enumerate(images):
        r, c = divmod(i, cols)
        left = grid_left + c * (cell_w + gap)
        top = grid_top + r * (cell_h + gap)
        path = prepare_image(item["image"], **item.get("image_args", {}))
        _add_image_fit(slide, path, left, top, cell_w, cell_h)

    text_left = grid_left + grid_w + Inches(0.4)
    _add_bullets(slide, text_left, Inches(1.9),
                 SLIDE_W - text_left - SIDE_M, Inches(4.6),
                 data["bullets"], size=17)
    _add_footer(slide, idx)


def render_triple_cards(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG)
    _add_header(slide, data["title"])
    top = Inches(1.7)
    if data.get("subtitle"):
        sub = slide.shapes.add_textbox(SIDE_M, Inches(1.55),
                                       SLIDE_W - 2 * SIDE_M, Inches(0.5))
        _set_text(sub.text_frame, data["subtitle"], font=FONT_BODY, size=15,
                  color=MUTED)
        top = Inches(2.15)
    _image_columns(slide, data["images"], top, Inches(3.5), card_bg=True)
    if data.get("note"):
        note = slide.shapes.add_textbox(SIDE_M, SLIDE_H - Inches(0.95),
                                        SLIDE_W - 2 * SIDE_M, Inches(0.5))
        _set_text(note.text_frame, data["note"], font=FONT_HEAD, size=16,
                  bold=True, color=NAVY)
    _add_footer(slide, idx)


def render_triple_media(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG)
    _add_header(slide, data["title"])
    _image_columns(slide, data["images"], Inches(1.9), Inches(3.9))
    if data.get("note"):
        note = slide.shapes.add_textbox(SIDE_M, SLIDE_H - Inches(0.9),
                                        SLIDE_W - 2 * SIDE_M, Inches(0.5))
        _set_text(note.text_frame, data["note"], font=FONT_BODY, size=12,
                  color=MUTED)
    _add_footer(slide, idx)


def render_factory_flow(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG)
    _add_header(slide, data["title"])

    left_w = Inches(5.6)
    if data.get("image"):
        path = prepare_image(data["image"], **data.get("image_args", {}))
        _add_image_fit(slide, path, SIDE_M, Inches(1.75), left_w, Inches(4.6))

    text_left = SIDE_M + left_w + Inches(0.4)
    text_w = SLIDE_W - text_left - SIDE_M
    _add_bullets(slide, text_left, Inches(1.85), text_w, Inches(2.4),
                 data["bullets"], size=17)
    if data.get("image2"):
        path2 = prepare_image(data["image2"], **data.get("image2_args", {}))
        _add_image_fit(slide, path2, text_left, Inches(4.35), text_w, Inches(2.0))
    _add_footer(slide, idx)


RENDERERS = {
    "title": render_title,
    "section": render_section,
    "split": render_split_text_image,
    "hero": render_image_hero,
    "dual": render_dual_image,
    "matrix": render_matrix,
    "title_branded": render_title_branded,
    "split_gallery": render_split_gallery,
    "triple_cards": render_triple_cards,
    "triple_media": render_triple_media,
    "factory_flow": render_factory_flow,
}

# ---------------------------------------------------------------------------
# Slide content (spec sections 5.1 / 5.2)
# ---------------------------------------------------------------------------

YUHUAN_GRID_CROP = (32, 106, 1478, 828)

# ---------------------------------------------------------------------------
# Branded additional slides (Jixi Huijie / CASC) from
# additonal pages/additional-slides-spec-ru.md
# ---------------------------------------------------------------------------

COMPANY_EN = "Jixi Huijie Aircraft Research and Development Center Co., Ltd."
COMPANY_CN = "鸡西市会杰飞行器研发中心有限责任公司"

ADDITIONAL_SLIDES: list[dict] = [
    {
        "layout": "title_branded",
        "title": COMPANY_EN,
        "subtitle": COMPANY_CN,
        "image": f"{VIDEO}/airplane.jpg",
        "image_args": {"top_mask": 0.08, "bottom_mask": 0.15,
                       "mask_color": (0, 120, 168)},
        "contacts": [
            "Директор: 耿会杰 (Гэн Хуэйцзе)   |   Тел.: +86 137 0463 1315",
            "E-mail: jixi_huijie@proton.me",
            "Адрес: ул. Авиационная — TBD",
        ],
    },
    {
        "layout": "split_gallery",
        "title": "Исследования, производство и международные поставки в авиации",
        "bullets": [
            "Научно-технические исследования и разработки в авиации.",
            "Производство и поставка легкомоторных самолётов, вертолётов "
            "и авиакомплектующих.",
            "Международные поставки и содействие импортно-экспортным операциям.",
        ],
        "images": [
            {"image": f"{VIDEO}/helocopter3.jpg",
             "image_args": {"top_mask": 0.12}},
            {"image": f"{VIDEO}/helocopter1.jpg",
             "image_args": {"top_mask": 0.12}},
            {"image": f"{VIDEO}/airplane.jpg",
             "image_args": {"top_mask": 0.08, "bottom_mask": 0.15}},
            {"image": f"{VIDEO}/agr_airplane.jpg",
             "image_args": {"bottom_mask": 0.28}},
        ],
    },
    {
        "layout": "triple_cards",
        "title": "Авиационные запчасти",
        "subtitle": "Деятельность по поставке запасных частей: таблицы, фото "
                    "и инвентарные номера.",
        "images": [
            {"image": f"{YUHUAN}/_xlsx_slide-03.png",
             "caption": "Каталог з/ч (инв. номера)"},
            {"image": f"{YUHUAN}/slide-12.png",
             "image_args": {"bottom_mask": 0.08, "mask_color": (255, 255, 255)},
             "caption": "Быстросъёмные замки"},
            {"image": f"{YUHUAN}/slide-13.png",
             "image_args": {"bottom_mask": 0.08, "mask_color": (255, 255, 255)},
             "caption": "Электрозамки"},
        ],
        "note": "Замки и запорные механизмы.",
    },
    {
        "layout": "triple_media",
        "title": "Фото и данные",
        "images": [
            {"image": f"{VIDEO}/helocopter2.jpg",
             "image_args": {"top_mask": 0.10},
             "caption": "Площадка и демонстрация техники"},
            {"image": f"{YUHUAN}/slide-10.png",
             "image_args": {"crop_px": YUHUAN_GRID_CROP},
             "caption": "Запасные части (обзор)"},
            {"image": f"{YUHUAN}/_xlsx_slide-02.png",
             "caption": "Данные каталога (Excel)"},
        ],
        "note": "Фото с аэродрома, изделия из презентации и таблицы каталога.",
    },
    {
        "layout": "split",
        "title": "CNAS-аккредитация",
        "bullets": [
            "Тестовые лаборатории.",
            "Аккредитация CNAS.",
            "Сертификаты качества: ISO 9001, отраслевые подтверждения.",
        ],
        "image": f"{VIDEO}/cert.png",
        "image_args": {"crop_frac": (0.05, 0.12, 0.90, 0.65),
                       "bottom_mask": 0.10, "mask_color": (255, 255, 255)},
        "caption": "Сертификаты и аккредитации.",
    },
    {
        "layout": "factory_flow",
        "title": "Работа с импортёрами",
        "bullets": [
            "Производство и поставка авиационных запчастей и комплектующих.",
            "Работа с импортёрами и международными партнёрами.",
            "Лицензированные перевозчики для авиационных грузов.",
        ],
        "image": f"{VIDEO}/production1.png",
        "image_args": {"bottom_mask": 0.14},
        "image2": f"{VIDEO}/production2.png",
        "image2_args": {"bottom_mask": 0.14},
    },
]

SLIDES: list[dict] = [
    {
        "layout": "title",
        "title": "Производственные возможности",
        "subtitle": "Авиационные системы запирания, GSE и крепёж высокого класса.\n"
                    "Обобщённый обзор без раскрытия поставщиков.",
    },
    {
        "layout": "matrix",
        "title": "Карта возможностей",
        "cards": [
            {
                "title": "A. Запирающие системы и GSE",
                "bullets": [
                    "Системы запирания дверей кабины и обтекателей",
                    "Быстросъёмные механизмы (quick-release)",
                    "Электрозамки, замки гондолы двигателя",
                    "Аварийный выпуск шасси",
                    "Наземное оборудование (GSE)",
                ],
            },
            {
                "title": "B. Крепёж высокого класса",
                "bullets": [
                    "Авиационный и космический крепёж",
                    "Железнодорожный крепёж и рельсовые костыли",
                    "Материалы: нерж., жаропрочные, титановые сплавы",
                    "Нестандартные изделия по чертежам",
                    "НИОКР, контроль качества, прослеживаемость",
                ],
            },
        ],
    },
    {
        "layout": "section",
        "letter": "A",
        "title": "Системы запирания и авиационные механизмы",
        "note": "Запирающие механизмы, quick-release, кабина, обтекатели, шасси и GSE.",
    },
    {
        "layout": "split",
        "title": "Авиационные системы запирания",
        "bullets": [
            "Системы запирания дверей кабины для вертолётов, eVTOL и БПЛА.",
            "Механизмы быстрой установки и снятия для обтекателей.",
            "Ориентация на требования летной годности.",
        ],
    },
    {
        "layout": "dual",
        "title": "Типовая линейка — общий обзор",
        "images": [
            {
                "image": f"{YUHUAN}/slide-10.png",
                "image_args": {"crop_px": YUHUAN_GRID_CROP},
                "caption": "Запирающие системы",
            },
            {
                "image": f"{YUHUAN}/slide-11.png",
                "image_args": {"crop_px": YUHUAN_GRID_CROP},
                "caption": "Комплекты и GSE",
            },
        ],
        "caption": "82 позиции представлены общими сетками: 52 позиции запирающих "
                   "систем и 30 позиций комплектов / GSE. Семейства: HB, TR-K, "
                   "TR-C, TR-D, TR-S, TR-B, TR-H, TR-Q.",
    },
    {
        "layout": "matrix",
        "title": "Семейства продукции",
        "header": ["Семейство", "Тип", "Как показать"],
        "col_widths": [0.22, 0.48, 0.30],
        "rows": [
            ["HB2", "Замки/защёлки авиационного стандарта HB", "1 строка, без артикулов"],
            ["HB6539-HB8231", "Стандартизированные авиационные компоненты", "1 строка"],
            ["TR-K101-K802", "Быстросъёмные замки разных типоразмеров", "1 строка"],
            ["TR-C", "Тросовые и монтажные комплекты", "1 строка"],
            ["TR-D", "Штифты, втулки, platform-компоненты", "1 строка"],
            ["TR-S", "Рукоятки управления", "1 строка"],
            ["TR-B / TR-H / TR-Q", "Стержни, крюки, штифты, блоки", "1 строка"],
            ["GSE", "Тележки, стойки, платформы, кронштейны", "1 строка"],
        ],
    },
    {
        "layout": "hero",
        "title": "Антивибрационные быстросъёмные замки",
        "image": f"{YUHUAN}/slide-12.png",
        "image_args": {"bottom_mask": 0.08, "mask_color": (255, 255, 255)},
        "caption": "Применение на обтекателях хвостового редуктора, хвостового "
                   "конуса, передней кромки, хвостовой балки.",
    },
    {
        "layout": "hero",
        "title": "Электрические замки",
        "image": f"{YUHUAN}/slide-13.png",
        "image_args": {"bottom_mask": 0.08, "mask_color": (255, 255, 255)},
        "caption": "Линейка электрических замков для авиационных систем.",
    },
    {
        "layout": "hero",
        "title": "Замки и распорки гондолы двигателя",
        "image": f"{YUHUAN}/slide-14.png",
        "image_args": {"bottom_mask": 0.08, "mask_color": (255, 255, 255)},
        "caption": "Замки и распорные элементы для двигательных гондол.",
    },
    {
        "layout": "hero",
        "title": "Двери кабины и запирающие механизмы",
        "image": f"{YUHUAN}/slide-15.png",
        "image_args": {"bottom_mask": 0.08, "mask_color": (255, 255, 255)},
        "caption": "Комплекты дверей кабины и механизмов запирания.",
    },
    {
        "layout": "split",
        "title": "Аварийный выпуск шасси",
        "bullets": [
            "Передняя и основная стойки.",
            "Работа при потере давления в гидросистеме.",
            "Резервный выпуск при отказе электроснабжения.",
        ],
        "image": f"{YUHUAN}/slide-16.png",
        "image_args": {"bottom_mask": 0.08, "mask_color": (255, 255, 255)},
    },
    {
        "layout": "matrix",
        "title": "Наземное оборудование (GSE)",
        "header": ["№", "Тип оборудования"],
        "col_widths": [0.08, 0.92],
        "rows": [
            ["1", "Подъёмно-транспортное оборудование"],
            ["2", "Буксировочная штанга"],
            ["3", "Устройства швартовки"],
            ["4", "Рабочие стремянки"],
            ["5", "Транспортные комплекты"],
            ["6", "Специализированная оснастка на заказ"],
        ],
    },
    {
        "layout": "matrix",
        "title": "Компетенции летной годности",
        "header": ["Аспект", "Содержание"],
        "col_widths": [0.32, 0.68],
        "rows": [
            ["Разработка", "Механизмы запирания дверей кабины"],
            ["Опыт по типам ВС", "AC313, AC311A, AC313A, GA20"],
            ["Сертификация", "Завершённые программы летной годности"],
        ],
    },
    {
        "layout": "section",
        "letter": "B",
        "title": "Крепёж и метизы высокого класса",
        "note": "Авиация, космос, железнодорожный транспорт, нестандартные изделия "
                "по чертежам.",
    },
    {
        "layout": "split",
        "title": "Профиль направления крепежа",
        "bullets": [
            "Полный цикл НИОКР, проектирования, производства и контроля.",
            "Стандартный и нестандартный крепёж.",
            "Материалы: углеродистые, легированные, нержавеющие, жаропрочные "
            "и титановые сплавы.",
        ],
    },
    {
        "layout": "hero",
        "title": "Авиакосмический крепёж (1)",
        "image": f"{QIFENG}/slide-11.png",
        "image_args": {"crop_frac": (0.06, 0.18, 0.88, 0.58), "bottom_mask": 0.12},
        "caption": "Винты стандарта HB и высокоблокирующие болты Lock Bolt.",
    },
    {
        "layout": "hero",
        "title": "Авиакосмический крепёж (2)",
        "image": f"{QIFENG}/slide-12.png",
        "image_args": {"crop_frac": (0.06, 0.18, 0.88, 0.58), "bottom_mask": 0.12},
        "caption": "Гидравлические детали и трубные соединения.",
    },
    {
        "layout": "hero",
        "title": "Железнодорожный подвижной состав",
        "image": f"{QIFENG}/slide-13.png",
        "image_args": {"crop_frac": (0.06, 0.20, 0.88, 0.55), "bottom_mask": 0.12},
        "caption": "Крепёж из нержавеющей стали для локомотивов и вагонов.",
    },
    {
        "layout": "hero",
        "title": "Инфраструктура пути",
        "image": f"{QIFENG}/slide-14.png",
        "image_args": {"crop_frac": (0.06, 0.20, 0.88, 0.55), "bottom_mask": 0.12},
        "caption": "Крепёж путевого хозяйства, рельсовых стыков и элементов "
                   "верхнего строения пути.",
    },
    {
        "layout": "split",
        "title": "Рельсовые костыли",
        "bullets": [
            "Резьба M20-M27.",
            "Длина до 300 мм.",
            "Применение на линиях национального значения.",
        ],
        "image": f"{QIFENG}/slide-15.png",
        "image_args": {"crop_frac": (0.08, 0.24, 0.84, 0.50), "bottom_mask": 0.12},
    },
    {
        "layout": "matrix",
        "title": "Сводная матрица возможностей",
        "header": ["Направление", "Ключевые компетенции", "Типовые изделия"],
        "col_widths": [0.26, 0.37, 0.37],
        "rows": [
            ["Запирание кабины и обтекателей",
             "Летная годность, 5-осевая обработка, контроль качества",
             "Быстросъёмные замки, электрозамки, двери кабины"],
            ["Системы шасси", "Опытная разработка, интеграция с типом ВС",
             "Аварийный выпуск шасси"],
            ["GSE", "Оснастка на заказ",
             "Подвеска, буксировка, швартовка, стремянки"],
            ["Авиационный крепёж", "HB, Lock Bolt, контроль геометрии",
             "Болты, винты, гидрофитинги, трубные соединения"],
            ["ЖД-крепёж", "Типоразмеры M20-M27, изделия для пути",
             "Рельсовые костыли, крепёж подвижного состава"],
            ["НИОКР и качество", "Патенты, CMM, MES/ERP, прослеживаемость",
             "Нестандарт по чертежам, Ti/жаропрочные сплавы"],
        ],
    },
    {
        "layout": "title",
        "title": "НИОКР, качество и следующий шаг",
        "bullets": [
            "Патентная база и собственные разработки.",
            "5-осевая обработка, координатно-измерительные системы (CMM).",
            "Испытания, MES/ERP, прослеживаемость партий.",
        ],
        "subtitle": "Подбор изделий и уточнение применимости — по техническому заданию.",
    },
]


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------


ALL_SLIDES: list[dict] = ADDITIONAL_SLIDES + SLIDES


def build(keep_temp: bool = False) -> None:
    assert len(ALL_SLIDES) == TOTAL_SLIDES, (
        f"expected {TOTAL_SLIDES} slides, got {len(ALL_SLIDES)}"
    )
    if TEMP_DIR.exists():
        shutil.rmtree(TEMP_DIR)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    used_images: list[str] = []
    for i, data in enumerate(ALL_SLIDES, start=1):
        renderer = RENDERERS[data["layout"]]
        slide = _new_blank(prs)
        renderer(slide, data, i)
        if data.get("image"):
            used_images.append(data["image"])
        if data.get("image2"):
            used_images.append(data["image2"])
        for item in data.get("images", []):
            used_images.append(item["image"])

    prs.save(OUTPUT_PPTX)

    print(f"[done] {OUTPUT_PPTX}")
    print(f"[info] slides: {len(prs.slides._sldIdLst)}")
    print(f"[info] images embedded: {len(used_images)}")
    for path in used_images:
        print(f"   - {path}")

    if not keep_temp and TEMP_DIR.exists():
        shutil.rmtree(TEMP_DIR)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--keep-temp", action="store_true",
                    help="keep cropped temp PNGs in _pptx_assets/")
    args = ap.parse_args()
    build(keep_temp=args.keep_temp)
    return 0


if __name__ == "__main__":
    sys.exit(main())
