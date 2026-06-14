#!/usr/bin/env python3
"""Build anonymized Zhengling steering-components PPTX (~16 slides).

Usage:
  cd aero_secondhend/ppt
  python build_zhengling_presentation_pptx.py
"""

from __future__ import annotations

import argparse
import shutil
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
MATERIALS_DIR = HERE / "maretials"
ZHENGLING_DIR = MATERIALS_DIR / "zhengling-machinery"
TEMP_DIR = ZHENGLING_DIR / "_pptx_assets"
OUTPUT_PPTX = ZHENGLING_DIR / "presentation-ru.pptx"

ZL = "zhengling-machinery/images"

WHITELIST = {
    f"{ZL}/slide-13.png",
    f"{ZL}/slide-14.png",
    f"{ZL}/slide-15.png",
    f"{ZL}/slide-16.png",
    f"{ZL}/slide-17.png",
    f"{ZL}/slide-18.png",
    f"{ZL}/slide-19.png",
    f"{ZL}/slide-20.png",
    f"{ZL}/slide-23.png",
    f"{ZL}/slide-24.png",
    f"{ZL}/slide-26.png",
    f"{ZL}/slide-27.png",
    f"{ZL}/slide-28.png",
    f"{ZL}/slide-32.png",
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
NAVY = RGBColor(0x0B, 0x2A, 0x4A)
ACCENT = RGBColor(0x00, 0x78, 0xA8)
BG = RGBColor(0xFF, 0xFF, 0xFF)
BG_SOFT = RGBColor(0xF5, 0xF7, 0xFA)
RULE = RGBColor(0xCB, 0xD5, 0xE1)
TEXT = RGBColor(0x1F, 0x2A, 0x44)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xF8, 0xFA, 0xFC)
ACCENT_SOFT = RGBColor(0xE8, 0xF4, 0xF8)
FONT_HEAD = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"
SIDE_M = Inches(0.55)
HEADER_H = Inches(0.50)
TOTAL_SLIDES = 16

MASK = {"bottom_mask": 0.14, "mask_color": (255, 255, 255)}
DESIGN_CROP = {"crop_frac": (0.0, 0.10, 1.0, 0.76), **MASK}
PROD_CROP = {"crop_frac": (0.05, 0.12, 0.90, 0.70), **MASK}


def _solid_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _set_bg(slide, color: RGBColor) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid_fill(rect, color)
    sp = rect._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def _set_text(tf, text, *, font=FONT_BODY, size=14, bold=False, color=TEXT,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP) -> None:
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _new_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_accent_line(slide, left, top, width, *, color=ACCENT, height=Pt(3)) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _solid_fill(bar, color)


def _add_header(slide, title: str) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
    _solid_fill(band, NAVY)
    title_box = slide.shapes.add_textbox(SIDE_M, Inches(0.62), SLIDE_W - 2 * SIDE_M, Inches(0.7))
    _set_text(title_box.text_frame, title, font=FONT_HEAD, size=26, bold=True, color=NAVY)
    _add_accent_line(slide, SIDE_M, Inches(1.32), Inches(2.4))


def _add_footer(slide, idx: int) -> None:
    box = slide.shapes.add_textbox(
        SLIDE_W - Inches(2.0), SLIDE_H - Inches(0.42), Inches(1.6), Inches(0.3)
    )
    _set_text(box.text_frame, f"{idx} / {TOTAL_SLIDES}", font=FONT_BODY, size=10,
              color=MUTED, align=PP_ALIGN.RIGHT)


def _add_bullets(slide, left, top, width, height, bullets, *, size=17, color=TEXT) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(size * 0.72)
        p.line_spacing = 1.15
        dot = p.add_run()
        dot.text = "•  "
        dot.font.name = FONT_BODY
        dot.font.size = Pt(size)
        dot.font.bold = True
        dot.font.color.rgb = ACCENT
        run = p.add_run()
        run.text = line
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color


def _add_table(slide, left, top, width, header, rows, *, col_widths=None,
               header_size=13, body_size=12) -> None:
    n_rows = len(rows) + 1
    n_cols = len(header)
    row_h = Inches(0.42)
    gfx = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_h * n_rows)
    table = gfx.table
    table.first_row = False
    table.horz_banding = False
    if col_widths:
        total = sum(col_widths)
        for c, frac in enumerate(col_widths):
            table.columns[c].width = Emu(int(width * frac / total))
    for c, label in enumerate(header):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.name = FONT_HEAD
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_ALT if r % 2 == 0 else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.name = FONT_BODY
            run.font.size = Pt(body_size)
            run.font.color.rgb = TEXT


def assert_allowed(rel_path: str) -> None:
    assert "catalog-items" not in rel_path, rel_path
    assert rel_path in WHITELIST, rel_path


def prepare_image(rel_path: str, *, crop_px=None, crop_frac=None,
                  bottom_mask: float = 0.0, mask_color=(255, 255, 255)) -> Path:
    assert_allowed(rel_path)
    src = MATERIALS_DIR / rel_path
    if not src.is_file():
        raise FileNotFoundError(src)
    img = Image.open(src).convert("RGB")
    w, h = img.size
    if crop_px is not None:
        img = img.crop(crop_px)
    elif crop_frac is not None:
        x, y, cw, ch = crop_frac
        img = img.crop((int(w * x), int(h * y), int(w * (x + cw)), int(h * (y + ch))))
    if bottom_mask > 0:
        cw, ch = img.size
        band_h = int(ch * bottom_mask)
        if band_h > 0:
            img.paste(Image.new("RGB", (cw, band_h), mask_color), (0, ch - band_h))
    TEMP_DIR.mkdir(parents=True, exist_ok=True)
    out = TEMP_DIR / f"crop__{rel_path.replace('/', '__')}"
    img.save(out, format="PNG")
    return out


def _add_image_fit(slide, img_path: Path, left, top, max_w, max_h) -> None:
    with Image.open(img_path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    draw_w, draw_h = int(iw * scale), int(ih * scale)
    left = left + (max_w - draw_w) // 2
    top = top + (max_h - draw_h) // 2
    pic = slide.shapes.add_picture(str(img_path), left, top, draw_w, draw_h)
    pic.line.color.rgb = RULE
    pic.line.width = Pt(1)


def render_title(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG_SOFT)
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), SLIDE_W - Inches(2.0), Inches(1.4))
    _set_text(title_box.text_frame, data["title"], font=FONT_HEAD, size=38, bold=True,
              color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_accent_line(slide, (SLIDE_W - Inches(2.2)) // 2, Inches(4.05), Inches(2.2))
    if data.get("subtitle"):
        sub = slide.shapes.add_textbox(Inches(1.5), Inches(4.35), SLIDE_W - Inches(3.0), Inches(1.6))
        _set_text(sub.text_frame, data["subtitle"], font=FONT_BODY, size=18, color=MUTED,
                  align=PP_ALIGN.CENTER)
    if data.get("bullets"):
        _add_bullets(slide, Inches(3.4), Inches(4.5), SLIDE_W - Inches(6.8), Inches(2.4),
                     data["bullets"], size=15)


def render_section(slide, data: dict, idx: int) -> None:
    _set_bg(slide, NAVY)
    letter = slide.shapes.add_textbox(SIDE_M, Inches(2.0), Inches(3.0), Inches(2.2))
    _set_text(letter.text_frame, data["letter"], font=FONT_HEAD, size=120, bold=True,
              color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    title = slide.shapes.add_textbox(Inches(3.0), Inches(2.9), SLIDE_W - Inches(3.6), Inches(1.4))
    _set_text(title.text_frame, data["title"], font=FONT_HEAD, size=34, bold=True,
              color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if data.get("note"):
        note = slide.shapes.add_textbox(Inches(3.05), Inches(4.2), SLIDE_W - Inches(3.6), Inches(1.0))
        _set_text(note.text_frame, data["note"], font=FONT_BODY, size=16, color=ACCENT_SOFT)


def render_split(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG)
    _add_header(slide, data["title"])
    _add_bullets(slide, SIDE_M, Inches(1.7), Inches(5.4), Inches(4.8), data["bullets"])
    if data.get("image"):
        path = prepare_image(data["image"], **data.get("image_args", {}))
        _add_image_fit(slide, path, Inches(6.4), Inches(1.7),
                       SLIDE_W - Inches(6.4) - SIDE_M, Inches(4.6))
    _add_footer(slide, idx)


def render_hero(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG)
    _add_header(slide, data["title"])
    path = prepare_image(data["image"], **data.get("image_args", {}))
    _add_image_fit(slide, path, SIDE_M, Inches(1.65), SLIDE_W - 2 * SIDE_M, Inches(4.55))
    if data.get("caption"):
        cap = slide.shapes.add_textbox(SIDE_M, Inches(6.45), SLIDE_W - 2 * SIDE_M, Inches(0.5))
        _set_text(cap.text_frame, data["caption"], font=FONT_BODY, size=12, color=MUTED,
                  align=PP_ALIGN.CENTER)
    _add_footer(slide, idx)


def render_dual(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG)
    _add_header(slide, data["title"])
    content_w = SLIDE_W - 2 * SIDE_M
    gap = Inches(0.4)
    cell_w = (content_w - gap) // 2
    cell_h = Inches(3.9)
    top = Inches(1.65)
    for i, item in enumerate(data["images"]):
        left = SIDE_M + i * (cell_w + gap)
        path = prepare_image(item["image"], **item.get("image_args", {}))
        _add_image_fit(slide, path, left, top, cell_w, cell_h)
        cap = slide.shapes.add_textbox(left, top + cell_h + Inches(0.05), cell_w, Inches(0.4))
        _set_text(cap.text_frame, item["caption"], font=FONT_HEAD, size=13, bold=True,
                  color=NAVY, align=PP_ALIGN.CENTER)
    _add_footer(slide, idx)


def render_matrix(slide, data: dict, idx: int) -> None:
    _set_bg(slide, BG)
    _add_header(slide, data["title"])
    if data.get("cards"):
        content_w = SLIDE_W - 2 * SIDE_M
        gap = Inches(0.35)
        n = len(data["cards"])
        cell_w = (content_w - gap * (n - 1)) // n
        top = Inches(1.75)
        cell_h = Inches(4.5)
        for i, card in enumerate(data["cards"]):
            left = SIDE_M + i * (cell_w + gap)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, cell_w, cell_h)
            box.fill.solid()
            box.fill.fore_color.rgb = BG_SOFT
            box.line.color.rgb = RULE
            head = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2),
                                            cell_w - Inches(0.4), Inches(0.6))
            _set_text(head.text_frame, card["title"], font=FONT_HEAD, size=16, bold=True, color=NAVY)
            _add_bullets(slide, left + Inches(0.2), top + Inches(0.95),
                         cell_w - Inches(0.4), cell_h - Inches(1.2),
                         card["bullets"], size=13)
    else:
        _add_table(slide, SIDE_M, Inches(1.7), SLIDE_W - 2 * SIDE_M,
                   data["header"], data["rows"], col_widths=data.get("col_widths"))
    _add_footer(slide, idx)


RENDERERS = {
    "title": render_title,
    "section": render_section,
    "split": render_split,
    "hero": render_hero,
    "dual": render_dual,
    "matrix": render_matrix,
}

SLIDES: list[dict] = [
    {
        "layout": "title",
        "title": "Компоненты рулевой системы",
        "subtitle": "Автомобильная отрасль: промежуточные валы, yoke, рулевые валы.\n"
                    "Обобщённый обзор без раскрытия поставщиков.",
    },
    {
        "layout": "matrix",
        "title": "Карта возможностей",
        "cards": [
            {"title": "Продукция", "bullets": ["I-Shaft 10T / 18T / EPS", "Yoke и рулевые валы", "OEM по образцам"]},
            {"title": "НИОКР", "bullets": ["28 технологов, 6 экспертов", "Стандарты EU/US/JP", "FEA и фазировка"]},
            {"title": "Производство", "bullets": ["3 линии сборки", "Сварка, spline, ball-sliding", "Холодная/горячая высадка"]},
            {"title": "Контроль", "bullets": ["Программа испытаний NVH", "CMM, MPT, стенды жёсткости", "100% контроль сырья"]},
        ],
    },
    {
        "layout": "section",
        "letter": "A",
        "title": "Продукция рулевой системы",
        "note": "Промежуточные валы, yoke, рулевые валы и карданные шарниры.",
    },
    {
        "layout": "split",
        "title": "НИОКР и проектирование",
        "bullets": [
            "Владение государственными и отраслевыми стандартами (EU, US, JP).",
            "Команда с опытом более 30 лет у ведущих производителей рулевых систем.",
            "Оптимизация конструкций по образцам мировых поставщиков компонентов.",
        ],
    },
    {
        "layout": "hero",
        "title": "Линейка промежуточных валов (I-Shaft)",
        "image": f"{ZL}/slide-14.png",
        "image_args": MASK,
        "caption": "НИОКР и серийное производство компонентов рулевой и трансмиссионной системы.",
    },
    {
        "layout": "hero",
        "title": "Yoke (вилки)",
        "image": f"{ZL}/slide-13.png",
        "image_args": MASK,
        "caption": "Холодная и горячая высадка, OEM по образцам заказчика.",
    },
    {
        "layout": "hero",
        "title": "10T скользящий I-Shaft",
        "image": f"{ZL}/slide-15.png",
        "image_args": MASK,
        "caption": "Сварной yoke, yoke холодной высадки, spline с пластиковым покрытием. "
                   "Рабочий момент: 30–50 N·m.",
    },
    {
        "layout": "hero",
        "title": "18T скользящий I-Shaft",
        "image": f"{ZL}/slide-16.png",
        "image_args": MASK,
        "caption": "Рабочий момент: 50–65 N·m.",
    },
    {
        "layout": "split",
        "title": "EPS ball-sliding I-Shaft",
        "bullets": [
            "Рабочий момент: 65–99 N·m.",
            "6-зубчатая конструкция со стальными шариками и втулками.",
            "Секционная сварка или монолитная сборка — по требованию заказчика.",
        ],
        "image": f"{ZL}/slide-17.png",
        "image_args": MASK,
    },
    {
        "layout": "hero",
        "title": "Рулевые валы",
        "image": f"{ZL}/slide-18.png",
        "image_args": MASK,
    },
    {
        "layout": "dual",
        "title": "Проектирование и FEA",
        "images": [
            {"image": f"{ZL}/slide-19.png", "image_args": DESIGN_CROP,
             "caption": "Карданный шарнир и фазировка"},
            {"image": f"{ZL}/slide-20.png", "image_args": DESIGN_CROP,
             "caption": "FEA-анализ деталей I-Shaft"},
        ],
    },
    {
        "layout": "matrix",
        "title": "Программа испытаний",
        "header": ["Категория", "Параметры"],
        "col_widths": [0.32, 0.68],
        "rows": [
            ["Функциональные", "Момент трения, соосность; крутильная жёсткость; "
             "осевой и радиальный зазоры; момент карданного шарнира"],
            ["Прочностные", "Предел кручения; усталость; усилие выдавливания подшипника; "
             "осевое растяжение вал–труба–yoke"],
            ["NVH", "Вязкое трение, акустические и вибрационные требования"],
        ],
    },
    {
        "layout": "section",
        "letter": "B",
        "title": "Производство и контроль",
        "note": "3 линии сборки, автоматическая сварка, стенды жёсткости и контроля качества.",
    },
    {
        "layout": "hero",
        "title": "Сборка промежуточного вала",
        "image": f"{ZL}/slide-28.png",
        "image_args": PROD_CROP,
    },
    {
        "layout": "matrix",
        "title": "Ключевые характеристики и материалы",
        "header": ["Аспект", "Содержание"],
        "col_widths": [0.35, 0.65],
        "rows": [
            ["Spline / подшипник", "Точность посадки spline; размеры отверстия подшипника; "
             "позиционирование болта"],
            ["Сырьё", "100% контроль сырья; оценка микроструктуры"],
            ["Процесс yoke", "Требования к штамповке и контролю дефектов материала"],
        ],
    },
    {
        "layout": "title",
        "title": "НИОКР, качество и следующий шаг",
        "subtitle": "Подбор изделий и уточнение применимости — по техническому заданию.",
        "bullets": [
            "FEA, оптимизация угла фазировки, программа испытаний NVH.",
            "3 линии сборки, сварка, контроль карданных шарниров.",
            "CMM, магнитопорошковый контроль, стенды жёсткости.",
        ],
    },
]

BANNED = ("Zhengling", "正菱", "永嘉", "Yongjia")


def build(keep_temp: bool = False) -> None:
    assert len(SLIDES) == TOTAL_SLIDES
    for data in SLIDES:
        blob = str(data)
        for word in BANNED:
            assert word not in blob, f"banned token {word!r} in slide data"

    if TEMP_DIR.exists():
        shutil.rmtree(TEMP_DIR)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    used: list[str] = []

    for i, data in enumerate(SLIDES, start=1):
        slide = _new_blank(prs)
        RENDERERS[data["layout"]](slide, data, i)
        if data.get("image"):
            used.append(data["image"])
        for item in data.get("images", []):
            used.append(item["image"])

    prs.save(OUTPUT_PPTX)
    print(f"[done] {OUTPUT_PPTX}")
    print(f"[info] slides: {len(prs.slides._sldIdLst)}")
    print(f"[info] images: {len(used)}")
    for p in used:
        print(f"   - {p}")

    if not keep_temp and TEMP_DIR.exists():
        shutil.rmtree(TEMP_DIR)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--keep-temp", action="store_true")
    args = ap.parse_args()
    build(keep_temp=args.keep_temp)
    return 0


if __name__ == "__main__":
    sys.exit(main())
